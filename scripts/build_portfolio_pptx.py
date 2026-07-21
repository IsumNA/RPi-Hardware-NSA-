#!/usr/bin/env python3
"""Build a large visual portfolio PowerPoint of NSA / IMX662 / Hailo results.

Gallery-style deck: dark frames, large images, clear captions (model, gain,
PSNR, params, TOPS/FPS/MAC where known). Regenerable from on-disk artifacts.
"""
from __future__ import annotations

import io
import json
import math
from pathlib import Path
from typing import Optional

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "portfolio" / "NSA_Results_Portfolio.pptx"

# Gallery palette
BG = RGBColor(0x12, 0x12, 0x14)
FG = RGBColor(0xF2, 0xF0, 0xEA)
MUTED = RGBColor(0xA8, 0xA4, 0x9A)
ACCENT = RGBColor(0xE8, 0xC9, 0x7A)
CARD = RGBColor(0x1C, 0x1C, 0x20)
RULE = RGBColor(0x3A, 0x38, 0x34)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Remove shadow if any
    sp = shape._element
    effect = sp.find(qn("p:effectLst"))
    if effect is not None:
        sp.remove(effect)


def _textbox(slide, left, top, width, height, text: str, *, size=14, bold=False,
             color=FG, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _multitext(slide, left, top, width, height, lines: list[tuple[str, dict]]):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, style) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = style.get("align", PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(style.get("size", 12))
        run.font.bold = style.get("bold", False)
        run.font.color.rgb = style.get("color", FG)
        run.font.name = style.get("font", "Calibri")
        if "space_after" in style:
            p.space_after = Pt(style["space_after"])
    return box


def _load_json(path: Path):
    if not path.exists():
        return None
    return json.loads(path.read_text())


def _pil_fit(path: Path, max_w: int = 1800, max_h: int = 1100, quality: int = 88) -> Optional[io.BytesIO]:
    if not path or not Path(path).exists():
        return None
    try:
        im = Image.open(path)
        im = ImageOps.exif_transpose(im)
        if im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        elif im.mode == "L":
            im = im.convert("RGB")
        im.thumbnail((max_w, max_h), Image.Resampling.LANCZOS)
        buf = io.BytesIO()
        im.save(buf, format="JPEG", quality=quality, optimize=True)
        buf.seek(0)
        return buf
    except Exception as e:
        print(f"  skip image {path}: {e}")
        return None


def _place_image(slide, path: Path, left, top, width, height) -> bool:
    buf = _pil_fit(path)
    if buf is None:
        return False
    # Fit inside box preserving aspect
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        iw, ih = 16, 9
    box_w = float(width)
    box_h = float(height)
    scale = min(box_w / iw, box_h / ih)
    dw = Emu(int(iw * scale))
    dh = Emu(int(ih * scale))
    # Center in box
    ox = left + Emu(int((box_w - float(dw)) / 2))
    oy = top + Emu(int((box_h - float(dh)) / 2))
    buf.seek(0)
    slide.shapes.add_picture(buf, ox, oy, width=dw, height=dh)
    return True


def section_slide(prs, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_rect(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.04), ACCENT)
    _textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.8),
             title, size=36, bold=True, color=FG, align=PP_ALIGN.LEFT)
    if subtitle:
        _textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.2),
                 subtitle, size=16, color=MUTED)


def title_slide(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), ACCENT)
    _textbox(slide, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.2),
             "Night-vision denoise results", size=40, bold=True, color=FG)
    _textbox(slide, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.8),
             "What this deck shows (in pictures):\n"
             "how well we can remove noise from a Sony IMX662 camera at high gain,\n"
             "and how fast that can run on a Raspberry Pi 5 + Hailo AI chip.",
             size=18, color=MUTED)
    _textbox(slide, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.2),
             "Every slide has a big plain-English title saying exactly what the image is.\n"
             "HCG = High Conversion Gain (the camera’s sensitive night mode). "
             "Gain 512 = darkest / noisiest setting we care about.",
             size=14, color=MUTED)


def caption_block(slide, left, top, width, height, title: str, body: str) -> None:
    """Bottom card: short what-it-is label + plain-English stats."""
    _add_rect(slide, left, top, width, height, CARD)
    lines = [
        (title, {"size": 12, "bold": True, "color": ACCENT, "space_after": 3}),
        (body, {"size": 11, "color": FG}),
    ]
    _multitext(slide, left + Inches(0.12), top + Inches(0.06),
               width - Inches(0.2), height - Inches(0.1), lines)


def image_plate(prs, path: Path, title: str, caption: str, *,
                badge: str = "", what: str = "What you are looking at") -> bool:
    """Full-bleed image with a big plain-English header."""
    if not path.exists():
        return False
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    # Big header — wraps to two lines; this is the outsider-facing sentence
    _textbox(slide, Inches(0.35), Inches(0.12), Inches(10.7), Inches(0.85),
             title, size=22, bold=True, color=FG)
    if badge:
        _textbox(slide, Inches(10.9), Inches(0.18), Inches(2.1), Inches(0.55),
                 badge, size=12, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    ok = _place_image(slide, path, Inches(0.35), Inches(1.0),
                      Inches(12.6), Inches(5.15))
    if not ok:
        return False
    caption_block(slide, Inches(0.35), Inches(6.25), Inches(12.6), Inches(1.05),
                  what, caption)
    return True


def two_up(prs, left_path: Path, right_path: Path, title: str,
           left_cap: str, right_cap: str, footer: str = "",
           left_what: str = "LEFT image", right_what: str = "RIGHT image") -> bool:
    if not left_path.exists() and not right_path.exists():
        return False
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _textbox(slide, Inches(0.35), Inches(0.1), Inches(12.6), Inches(0.7),
             title, size=20, bold=True, color=FG)
    if left_path.exists():
        _place_image(slide, left_path, Inches(0.3), Inches(0.85),
                     Inches(6.3), Inches(4.55))
        caption_block(slide, Inches(0.3), Inches(5.5), Inches(6.3), Inches(1.2),
                      left_what, left_cap)
    if right_path.exists():
        _place_image(slide, right_path, Inches(6.75), Inches(0.85),
                     Inches(6.3), Inches(4.55))
        caption_block(slide, Inches(6.75), Inches(5.5), Inches(6.3), Inches(1.2),
                      right_what, right_cap)
    if footer:
        _textbox(slide, Inches(0.35), Inches(6.85), Inches(12.6), Inches(0.45),
                 footer, size=11, color=MUTED)
    return True


def four_grid(prs, paths: list[Path], title: str, captions: list[str]) -> bool:
    paths = [p for p in paths if p.exists()][:4]
    if not paths:
        return False
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _textbox(slide, Inches(0.3), Inches(0.08), Inches(12.7), Inches(0.55),
             title, size=18, bold=True, color=FG)
    positions = [
        (0.3, 0.7), (6.85, 0.7),
        (0.3, 4.0), (6.85, 4.0),
    ]
    for i, p in enumerate(paths):
        x, y = positions[i]
        _place_image(slide, p, Inches(x), Inches(y), Inches(6.15), Inches(2.7))
        cap = captions[i] if i < len(captions) else p.name
        _textbox(slide, Inches(x), Inches(y + 2.7), Inches(6.15), Inches(0.45),
                 cap, size=11, bold=True, color=ACCENT)
    return True


def bullets_slide(prs, title: str, rows: list[str], subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
             title, size=28, bold=True, color=FG)
    if subtitle:
        _textbox(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.5),
                 subtitle, size=14, color=MUTED)
    y = 1.7 if subtitle else 1.25
    lines = [(f"·  {r}", {"size": 16, "color": FG, "space_after": 10}) for r in rows]
    _multitext(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(5.3), lines)


def metrics_table_slide(prs, title: str, headers: list[str], rows: list[list[str]],
                        note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _textbox(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
             title, size=22, bold=True, color=FG)
    table = slide.shapes.add_table(len(rows) + 1, len(headers),
                                   Inches(0.5), Inches(1.1),
                                   Inches(12.3), Inches(0.45 * (len(rows) + 1))).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)
                r.font.color.rgb = ACCENT
                r.font.name = "Calibri"
        # dark cell fill
        _shade_cell(cell, "1C1C20")
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.color.rgb = FG
                    r.font.name = "Calibri"
            _shade_cell(cell, "121214" if i % 2 == 0 else "18181C")
    if note:
        _textbox(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.45),
                 note, size=11, color=MUTED)


def _shade_cell(cell, hex6: str) -> None:
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solid = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{hex6}"/></a:solidFill>'
    )
    # remove existing solidFill
    for child in list(tcPr):
        if "solidFill" in child.tag:
            tcPr.remove(child)
    tcPr.append(solid)


# ---------------------------------------------------------------------------
# Content builders
# ---------------------------------------------------------------------------

def build_overview(prs) -> None:
    title_slide(prs)
    bullets_slide(
        prs,
        "What is in this deck (in order)",
        [
            "1. How fast each chip can run the denoiser (Hailo vs Pi CPU vs NPU)",
            "2. Proof our fake camera noise matches the real IMX662 (especially HCG gain 512)",
            "3. The synthetic training pictures we generated from that noise model",
            "4. Before / after denoise results from our best neural models",
            "5. Teacher vs fast student model comparisons at gain 128 and 512",
            "6. Live-camera and training progress pictures",
            "7. Older architecture-search experiments (for context)",
        ],
        subtitle="You do not need the jargon — each image slide states what it is in the big title.",
    )
    bullets_slide(
        prs,
        "The setup in one page",
        [
            "Camera: Sony IMX662 — night modes LCG (normal) and HCG (more sensitive)",
            "Gains we test: 128, 256, 512 — higher gain = darker scene, more noise",
            "Computer: Raspberry Pi 5 with a Hailo-10H AI chip (40 TOPS INT8)",
            "Goal: clean up the noisy camera image in real time (≥30 frames per second)",
            "Best small model so far: ~196,000 weights, trained on gains 128+256+512",
            "PSNR in dB = quality score (higher is better; ~28 dB is our current holdout range)",
        ],
    )


def build_compute(prs) -> None:
    section_slide(
        prs,
        "Part 1 — How fast can it run?",
        "Chip speed (TOPS), frames per second, and how cost grows with image size",
    )
    metrics_table_slide(
        prs,
        "Chip speed at a glance — who can hit real-time video?",
        ["Chip", "Peak TOPS", "Number format", "Typical FPS", "Time per frame", "Plain meaning"],
        [
            ["Hailo-10H (on Pi)", "40", "INT8", "≥30 goal", "<33 ms goal", "Main target for live video"],
            ["Hailo-8 (older test)", "26", "INT8", "50.8", "19.7 ms", "Past run — already real-time"],
            ["DeepX DX-M1", "25", "INT8", "41–47", "~22 ms", "Alternative AI chip"],
            ["Pi 5 CPU alone", "0.04", "FP16", "~6–10", "100–164 ms", "Too slow for smooth video"],
            ["Intel NPU (server)", "—", "FP16", "182", "5.5 ms", "Fast on a PC, not on the Pi"],
        ],
        note="TOPS = trillions of AI ops per second. FPS = frames cleaned per second. Hailo-on-Pi is the real product path.",
    )
    metrics_table_slide(
        prs,
        "Same denoiser, different image sizes — CPU vs NPU on a PC (not the Pi)",
        ["Image size", "CPU time", "CPU FPS", "NPU time", "NPU FPS"],
        [
            ["Large 968×544", "1262 ms", "0.8", "217 ms", "4.6"],
            ["Medium 484×272", "331 ms", "3.0", "50 ms", "19.8"],
            ["Small 240×144", "62 ms", "16.0", "14 ms", "69.6"],
        ],
        note="Bigger images = slower. Sending full frames over WiFi from Pi→PC is also too slow (~1.5 FPS) — so we run Hailo on the Pi itself.",
    )

    hw_plain = [
        ("20260708-223649_nafnet_hailo8",
         "Chart: how much AI compute (TOPS) this model needs as the image gets bigger — Hailo-8 test",
         "HAILO"),
        ("20260714-210442_nafnet_rpi5_cpu",
         "Chart: how much AI compute (TOPS) this model needs as the image gets bigger — Pi 5 CPU",
         "PI CPU"),
        ("20260711-141214_nafnet_intel_npu",
         "Chart: how much AI compute (TOPS) this model needs as the image gets bigger — Intel NPU",
         "INTEL NPU"),
    ]
    for hw_dir, title, badge in hw_plain:
        p = ROOT / "outputs" / "history" / hw_dir / "resolution_tops_scaling.png"
        summ = _load_json(ROOT / "outputs" / "history" / hw_dir / "summary.json") or {}
        model = summ.get("model", {})
        params = model.get("params")
        params_s = f"{params:,}" if isinstance(params, int) else str(params)
        cap = (
            f"Hardware tested: {summ.get('hardware_name', hw_dir)}. "
            f"Model: {model.get('family', '?')} with {params_s} weights.\n"
            f"Speed in this run: about {summ.get('fps')} FPS "
            f"({summ.get('latency_ms')} ms per frame), {summ.get('precision')} math.\n"
            f"Image quality: PSNR {summ.get('psnr_in')} → {summ.get('psnr_out')} dB "
            f"(noisy → cleaned). Dashed line = chip’s peak TOPS."
        )
        image_plate(prs, p, title, cap, badge=badge, what="This chart means")

    image_plate(
        prs,
        ROOT / "outputs" / "resolution_tops_scaling.png",
        "Chart: one plot comparing Pi CPU, Hailo, DeepX, and Intel NPU compute vs image size",
        "As the camera image gets larger, the model needs more multiply-add work (MAC / GFLOPs). "
        "Pi CPU maxes out near 0.04 TOPS; Hailo-class chips have much more headroom for live video.",
        badge="ALL CHIPS",
        what="This chart means",
    )

    for hw_dir, badge, plain_chip in [
        ("20260708-223649_nafnet_hailo8", "HAILO", "Hailo AI chip"),
        ("20260714-210442_nafnet_rpi5_cpu", "PI CPU", "Raspberry Pi 5 CPU only"),
        ("20260711-141214_nafnet_intel_npu", "INTEL NPU", "Intel PC NPU"),
    ]:
        panel = ROOT / "outputs" / "history" / hw_dir / "validation_panel.png"
        summ = _load_json(ROOT / "outputs" / "history" / hw_dir / "summary.json") or {}
        model = summ.get("model", {})
        cap = (
            f"Three panels: noisy camera crop → model output → clean ground truth. "
            f"Run on {plain_chip}.\n"
            f"Model family {model.get('family', '?')}, {model.get('params', '?')} weights, "
            f"sensor {summ.get('sensor')} at gain {summ.get('gain')}.\n"
            f"Quality: PSNR {summ.get('psnr_in')} → {summ.get('psnr_out')} dB · "
            f"Speed: {summ.get('fps')} FPS ({summ.get('latency_ms')} ms)."
        )
        image_plate(
            prs, panel,
            f"Before / after denoise example — this run was on the {plain_chip}",
            cap, badge=badge, what="What the three panels are",
        )


def build_noise(prs) -> None:
    section_slide(
        prs,
        "Part 2 — Does our fake noise look like the real camera?",
        "We measured the IMX662 in HCG mode, then checked synthetic noise against real noise",
    )
    summ = _load_json(ROOT / "datasets" / "noise_profiles" / "hcg_validation" / "summary.json") or {}
    gains = summ.get("gains", {})
    verdict = summ.get("verdict", {})

    rows = []
    for g in ("128", "256", "512"):
        d = gains.get(g, {})
        rows.append([
            f"HCG gain {g}",
            f"{d.get('sigma_err_rel_mean', 0)*100:.1f}%",
            f"{d.get('kl_div_mean', 0):.3f}",
            f"{d.get('model_ptc_r2_mean', 0):.3f}",
            "GOOD ENOUGH" if d.get("eval_ok") else "close / partial",
            "real camera burst or lab calib",
        ])
    metrics_table_slide(
        prs,
        "How close is fake noise to real noise? (lower error = better match)",
        ["Camera setting", "Noise-level error", "Shape mismatch (KL)", "Curve fit R²", "Pass?", "Data used"],
        rows,
        note=(verdict.get("practical") or "")[:240]
        or "Gain 512 HCG is the hardest night case and now passes our training gate.",
    )

    bullets_slide(
        prs,
        "How to read the next noise pictures",
        [
            "Side-by-side strip = REAL camera grain | OUR fake grain | difference (should look quiet)",
            "Photon-transfer chart = how noise grows as the image gets brighter (real measured data)",
            "Model overlay chart = our mathematical fit on top of that measured data",
            "Score panel = pass/fail numbers we use before generating training images",
            "HCG gain 512 is the star: ~4% noise-level error after the hybrid fix — good for training",
            "If fake grain does not match real grain, the neural net learns the wrong problem",
        ],
    )

    vdir = ROOT / "datasets" / "noise_profiles" / "hcg_validation"
    pdir = ROOT / "datasets" / "noise_profiles"

    for g, tag in [(128, "imx662h_ag128"), (256, "imx662h_ag256"), (512, "imx662h_ag512")]:
        d = gains.get(str(g), {})
        badge = f"HCG {g}"
        err = d.get("sigma_err_rel_mean", 0) * 100
        kl = d.get("kl_div_mean", 0)
        r2 = d.get("model_ptc_r2_mean", 0)
        stats_short = (
            f"IMX662 HCG mode, analog gain {g}. "
            f"Noise-level error {err:.1f}% · shape mismatch KL={kl:.3f} · curve R²={r2:.3f}."
        )
        two_up(
            prs,
            vdir / f"{tag}_sidebyside_strip.jpg",
            vdir / f"{tag}_sidebyside.png",
            f"Real camera grain vs our fake grain — HCG gain {g}",
            f"{stats_short} Strip order: real | fake | difference.",
            f"Same comparison as a larger panel. "
            f"Residual noise strength: real={d.get('visual', {}).get('sigma_resid_real', 0):.4f}, "
            f"fake={d.get('visual', {}).get('sigma_resid_synth', 0):.4f}.",
            footer="If the middle (fake) looks like the left (real), we can trust synthetic training data at this gain.",
            left_what="LEFT — strip: real | fake | difference",
            right_what="RIGHT — full side-by-side panel",
        )
        two_up(
            prs,
            vdir / f"{tag}_ptc_profile.png",
            vdir / f"{tag}_ptc_model.png",
            f"Noise-vs-brightness curves measured on the real camera — HCG gain {g}",
            "Measured photon-transfer data from real captures (one curve per colour channel).",
            "Our fitted math model on top of that data (shot noise + read noise).",
            footer=stats_short,
            left_what="LEFT — measured from real camera",
            right_what="RIGHT — our fitted noise model",
        )
        image_plate(
            prs,
            vdir / f"{tag}_eval_panel.png",
            f"Scorecard: is HCG gain {g} fake noise good enough to train on?",
            f"{stats_short}\n"
            f"Pass rules we use: noise-level error under ~25% and KL under ~0.20. "
            f"Status: {'PASS' if d.get('eval_ok') else 'partial — see numbers'}.",
            badge=badge,
            what="What this scorecard shows",
        )
        image_plate(
            prs,
            pdir / f"imx662_ag{g}_ptc.png",
            f"Noise curve for the OTHER camera mode (LCG) at gain {g} — for comparison",
            f"LCG = Low Conversion Gain (less sensitive than HCG). "
            f"About 15% of training images use LCG so the model still works outside night-HCG.",
            badge="LCG",
            what="What this chart is",
        )
        image_plate(
            prs,
            pdir / f"imx662h_ag{g}_ptc.png",
            f"Official HCG gain {g} noise profile we use to generate fake training images",
            f"This is the stored IMX662 HCG×{g} profile. "
            f"Synthetic datasets draw Poisson/read/row noise from this file.",
            badge="HCG",
            what="What this chart is",
        )

    ba = verdict.get("ag512_before_after", {})
    if ba:
        bullets_slide(
            prs,
            "Fixing HCG gain 512 noise — before vs after (plain English)",
            [
                f"Before the fix: noise error {ba['before']['sigma_err']*100:.1f}%, "
                f"shape mismatch {ba['before']['kl']:.2f} — not good enough for training",
                f"After the fix: noise error {ba['after']['sigma_err']*100:.1f}%, "
                f"shape mismatch {ba['after']['kl']:.2f} — passes the gate",
                "What we changed: blend real burst measurements with lab calibration "
                "(especially the blue channel)",
                "Why it matters: gain 512 HCG is the dark livestream case we care about most",
            ],
        )


def _synth_pair_title(m: dict) -> str:
    kind = m.get("kind", "synth")
    gain = m.get("gain", "?")
    mode = m.get("mode", "HCG")
    pair = m.get("pair_kind", "synth")
    if pair == "real_burst":
        return f"Real camera noisy frame next to its clean average — {mode} gain {gain}"
    if kind == "live_like":
        return f"Synthetic noisy image made like the live dark stream — {mode} gain {gain}"
    if kind == "milder":
        return f"Synthetic noisy image at medium brightness — {mode} gain {gain}"
    return f"Synthetic noisy image generated with {mode} gain {gain} noise — next to clean target"


def build_synth(prs) -> None:
    section_slide(
        prs,
        "Part 3 — Fake training pictures we generated",
        "12,000 noisy/clean pairs, mostly HCG, including a slice that looks like the live dark stream",
    )
    bullets_slide(
        prs,
        "What is in the synthetic dataset (no jargon)",
        [
            "12,000 training pairs: each pair = one noisy crop + one clean crop (256×256)",
            "About 85% are HCG night mode; 15% are LCG so the model is not night-only",
            "Nearly half are the hardest setting: HCG gain 512",
            "Some pairs are deliberately as dark as the real Pi livestream (live_dark)",
            "Noise is painted on using the measured IMX662 profiles from Part 2",
            "Clean pictures come from real burst averages and other clean sources",
        ],
    )

    ex_dir = ROOT / "datasets" / "synth_pairs_hcg" / "examples"
    meta = _load_json(ex_dir / "examples_meta.json") or []

    sbs = sorted(ex_dir.glob("*_sbs.png"))
    four_grid(
        prs,
        sbs[:4],
        "Four training examples — each tile is noisy image | clean target",
        [
            "HCG 512 synth pair",
            "HCG 512 synth pair",
            "Live-dark style HCG 512",
            "Milder / mid-brightness pair",
        ][: len(sbs[:4])],
    )
    if len(sbs) > 4:
        caps = []
        for p in sbs[4:8]:
            name = p.stem.replace("_sbs", "")
            if "live_like" in name:
                caps.append("Looks like live dark stream · HCG 512")
            elif "milder" in name and "256" in name:
                caps.append("Milder scene · HCG gain 256")
            elif "milder" in name:
                caps.append("Milder scene · HCG gain 512")
            elif "hcg512" in name:
                caps.append("Synthetic pair · HCG gain 512")
            else:
                caps.append(name)
        four_grid(
            prs,
            sbs[4:8],
            "More training examples — noisy | clean side by side",
            caps,
        )

    for m in meta:
        png = Path(m.get("png") or "")
        if not png.exists():
            jpg = m.get("jpg") or ""
            png = ex_dir / Path(jpg).name.replace(".jpg", ".png")
        pair = m.get("pair_kind", "synth")
        origin = (
            "Painted with our HCG noise model on a clean source image"
            if pair == "synth"
            else "Taken from a real camera burst (not painted)"
        )
        cap = (
            f"{origin}. Mode {m.get('mode')} · gain {m.get('gain')} · "
            f"brightness band “{m.get('intensity_band')}”.\n"
            f"Source tag: {m.get('source')}. "
            f"Average brightness {m.get('gt_mean'):.3f} · "
            f"noise strength {m.get('resid_sigma'):.3f} "
            f"(display boosted ×{m.get('display_gain'):.1f} so you can see it)."
        )
        image_plate(
            prs, png, _synth_pair_title(m), cap,
            badge=f"GAIN {m.get('gain')}",
            what="What this side-by-side is",
        )

    for i, stem in enumerate([
        "hcg512_00_000000_0836",
        "hcg512_01_000002_0801",
        "hcg512_02_000003_0883",
    ], start=1):
        two_up(
            prs,
            ex_dir / f"{stem}_noisy.png",
            ex_dir / f"{stem}_gt.png",
            f"Example {i}: synthetic noisy crop (HCG gain 512) and its clean ground truth",
            "Fake noisy image — generated with the HCG gain 512 noise profile (brightened for display).",
            "Clean target the model should recover — same crop, no painted noise.",
            footer="This is what one training pair looks like when split apart.",
            left_what="LEFT — synthetic noisy image (HCG 512)",
            right_what="RIGHT — clean ground truth",
        )
        image_plate(
            prs,
            ex_dir / f"{stem}_canvas.jpg",
            f"Same HCG 512 training pair as a labeled gallery card (example {i})",
            "Website-style card: shows the pair with clear labels for noisy vs clean.",
            badge="HCG 512",
            what="What this card shows",
        )

    live = ROOT / "datasets" / "synth_pairs_hcg" / "live_ref"
    image_plate(
        prs,
        live / "live_sbs.jpg",
        "Real Pi livestream crop next to a clean burst average — HCG gain 512 reference",
        "This is the real dark live look we matched when building the live_dark synthetic slice.",
        badge="REAL LIVE",
        what="What you are looking at",
    )
    image_plate(
        prs,
        live / "user_screenshot_left_noisy.png",
        "Screenshot from the real livestream — the noisy left pane at HCG gain 512",
        "This is the worst-case live dark appearance. Synthetic live_dark pairs are built to look like this.",
        badge="REAL STREAM",
        what="What you are looking at",
    )
    image_plate(
        prs,
        ex_dir / "user_screenshot_left_crop256.png",
        "Same livestream noise, cropped to the 256×256 size used in training",
        "Cropped so you can compare grain size/strength directly against synthetic 256×256 pairs.",
        badge="REAL STREAM",
        what="What you are looking at",
    )


def _champion_caption(run_name: str, summ: dict) -> str:
    evals = summ.get("eval") or []
    eval_bits = []
    for e in evals[:3]:
        eval_bits.append(
            f"{e.get('scene')} at gain {e.get('gain')}: "
            f"{e.get('psnr_in')}→{e.get('psnr_out')} dB"
        )
    loss = summ.get("sample_loss", "?")
    params = summ.get("params", "?")
    gains = summ.get("gains", "?")
    return (
        f"Model run “{run_name}”: small 1-step denoiser with {params} weights, "
        f"trained with {loss} loss on gains {gains}.\n"
        f"Average quality: PSNR {summ.get('psnr_in')} → {summ.get('psnr_out')} dB "
        f"(noisy → cleaned). Sharpness ratio ≈ {summ.get('grad_ratio')}. "
        f"Training pairs: {summ.get('pairs') or summ.get('total_pairs')}.\n"
        + ("Holdout checks: " + " · ".join(eval_bits) if eval_bits else
           "Holdout scores are in the run summary.")
    )


def build_champions(prs) -> None:
    section_slide(
        prs,
        "Part 4 — Best denoise models (before / after pictures)",
        "Small fast neural nets (~196k weights) that clean noisy IMX662 frames in one step",
    )
    metrics_table_slide(
        prs,
        "Named best models — quality score after cleaning (higher PSNR is better)",
        ["Model name", "PSNR after", "Weights", "Training loss", "Gains trained", "In plain words"],
        [
            ["r101 dump", "28.39 dB", "196,292", "L1 + grad", "128–512", "Early strong student"],
            ["r128 balanced", "28.46 dB", "196,292", "L1", "128–512", "Balanced training"],
            ["r133 lpips", "28.39 dB", "196,292", "L1 + LPIPS", "128–512", "Looks better to eyes"],
            ["r134 hf-heavy", "28.49 dB", "196,292", "L1", "128–512", "More high-frequency detail"],
            ["r142 lpips long", "28.39 dB", "196,292", "L1 + LPIPS", "128–512", "Long perceptual train"],
            ["r209 / r221", "28.42 dB", "196,292", "L1 + LPIPS", "128–512", "Latest balanced line"],
        ],
        note="All of these are fast 1-step students meant for Hailo. Panels below show noisy → cleaned → true clean.",
    )

    viz = ROOT / "outputs" / "claude_champion_r142_viz"
    r142 = _load_json(ROOT / "cloud_pack" / "checkpoints" / "champion_r142_lpips_long_summary.json") or {}
    image_plate(
        prs,
        viz / "panel_noisy_denoise_gt.png",
        "Champion model r142 — noisy input | model output | true clean image",
        _champion_caption("r142_lpips_long", r142),
        badge="BEST r142",
        what="Left → middle → right",
    )
    two_up(
        prs,
        viz / "training_panel_r142.png",
        viz / "prior_champion_r134_panel.png",
        "Comparing two champions: r142 (newer) vs r134 (previous best)",
        _champion_caption("r142_lpips_long", r142),
        "Previous champion r134 — same style panel for a visual side-by-side of quality.",
        left_what="LEFT — model r142 result panel",
        right_what="RIGHT — older model r134 result panel",
    )
    two_up(
        prs,
        viz / "full_cabinet_H_2_ag128_f437.png",
        viz / "full_cabinet_D50_100_ag128_f429.png",
        "Full real scenes cleaned by r142 at gain 128 (two different rooms / lights)",
        "Scene “cabinet_H_2”, gain 128, holdout frame 437 — noisy vs cleaned vs clean reference.",
        "Scene “cabinet_D50_100”, gain 128, holdout frame 429 — same layout.",
        left_what="LEFT — scene cabinet_H_2 @ gain 128",
        right_what="RIGHT — scene cabinet_D50_100 @ gain 128",
    )
    four_grid(
        prs,
        [
            viz / "crops" / "crop_cabinet_H_2_ag128_f437.png",
            viz / "crops" / "crop2_cabinet_H_2_ag128_f437.png",
            viz / "crops" / "crop_cabinet_D50_100_ag128_f429.png",
            viz / "crops" / "crop2_cabinet_D50_100_ag128_f429.png",
        ],
        "Zoomed crops from r142 — check edges and texture (not just overall blur)",
        [
            "Crop: cabinet_H_2 detail A",
            "Crop: cabinet_H_2 detail B",
            "Crop: D50_100 detail A",
            "Crop: D50_100 detail B",
        ],
    )
    image_plate(
        prs,
        viz / "flicker_cabinet_H_2_ag128.png",
        "Several frames in a row from r142 — checking that the clean video does not flicker",
        "Same scene at gain 128 across time. We want stable brightness/texture, not hopping grain.",
        badge="TIME",
        what="What this strip shows",
    )

    named = [
        "r101_dump", "r128_cfm_bal", "r134_cfm_bal", "r142_lpips_long",
        "r147_cfm_bal", "r148_cfm_bal", "r164_cfm_bal", "r172_cfm_bal",
        "r196_cfm_bal", "r200_cfm_bal", "r209_cfm_bal", "r217_cfm_bal", "r221_cfm_bal",
    ]
    perfect = ROOT / "outputs" / "perfect_run"
    for name in named:
        cands = sorted(perfect.glob(f"*{name}*"))
        if not cands:
            cands = [perfect / name] if (perfect / name).exists() else []
        for d in cands[:1]:
            panel = d / "cfm_student_panel.png"
            summ = _load_json(d / "cfm_student_summary.json") or {}
            if not summ:
                for cp in (ROOT / "cloud_pack" / "checkpoints").glob(f"*{name}*summary.json"):
                    summ = _load_json(cp) or {}
            psnr = summ.get("psnr_out", "?")
            image_plate(
                prs,
                panel,
                f"Denoise result panel — model {d.name} (about {psnr} dB after cleaning)",
                _champion_caption(d.name, summ),
                badge=f"{psnr} dB",
                what="Usually: noisy | cleaned | true clean",
            )

    scored = []
    unscored = []
    for d in perfect.iterdir():
        if not d.is_dir():
            continue
        summ = _load_json(d / "cfm_student_summary.json") or {}
        panel = d / "cfm_student_panel.png"
        if not panel.exists():
            continue
        psnr = summ.get("psnr_out")
        if psnr is None:
            unscored.append((d, summ))
        else:
            scored.append((float(psnr), d, summ))
    scored.sort(reverse=True)
    shown = set(named)
    extra = [t for t in scored if not any(n in t[1].name for n in shown)]
    npages = math.ceil(len(extra) / 4) or 1
    for i in range(0, len(extra), 4):
        chunk = extra[i:i + 4]
        four_grid(
            prs,
            [t[1] / "cfm_student_panel.png" for t in chunk],
            f"More trained models ranked by clean quality (page {i//4 + 1} of {npages})",
            [
                f"{t[1].name}: {t[0]:.2f} dB after clean ({t[2].get('sample_loss', '?')} loss)"
                for t in chunk
            ],
        )
    for psnr, d, summ in extra[:16]:
        image_plate(
            prs,
            d / "cfm_student_panel.png",
            f"Close-up result — model {d.name} reached {psnr:.2f} dB after cleaning",
            _champion_caption(d.name, summ),
            badge=f"{psnr:.2f} dB",
            what="Usually: noisy | cleaned | true clean",
        )
    # Panels with a picture but no saved PSNR — keep, but say score missing
    if unscored:
        for i in range(0, min(len(unscored), 24), 4):
            chunk = unscored[i:i + 4]
            four_grid(
                prs,
                [t[0] / "cfm_student_panel.png" for t in chunk],
                "Extra result panels — quality score was not saved for these runs",
                [f"{t[0].name}: picture only (no PSNR in summary)" for t in chunk],
            )


def _humanize_edm_crop(stem: str) -> str:
    s = stem
    bits = []
    if "ag512" in s:
        bits.append("gain 512")
    elif "ag128" in s:
        bits.append("gain 128")
    if "teacher" in s:
        bits.append("slow teacher model")
    elif "student" in s:
        bits.append("fast student model")
    if "edm" in s:
        bits.append("EDM")
    if "l1" in s:
        bits.append("L1 train")
    if "l2" in s:
        bits.append("L2 train")
    if "euler" in s:
        bits.append("Euler sampler")
    if "heun" in s:
        bits.append("Heun sampler")
    if "charb" in s:
        bits.append("Charbonnier+edge")
    return " · ".join(bits) if bits else stem


def build_edm(prs) -> None:
    section_slide(
        prs,
        "Part 5 — Slow high-quality teacher vs fast student",
        "Teachers take many steps; students copy them in one step for real-time Hailo",
    )
    edm = ROOT / "outputs" / "edm_compare"
    image_plate(
        prs,
        edm / "teacher_grid_ag512.png",
        "Grid of slow teacher denoisers on a hard gain-512 scene (cabinet_H_2)",
        "Each tile is a different multi-step teacher. These are quality references, too slow for live video.",
        badge="TEACHERS",
        what="What this grid is",
    )
    two_up(
        prs,
        edm / "student_grid_ag128.png",
        edm / "student_grid_ag512.png",
        "Fast one-step students — easier gain 128 (left) vs harder gain 512 (right)",
        "Same student recipes at gain 128 — less noise, easier clean-up.",
        "Same student recipes at gain 512 — much noisier; this is the stress test.",
        left_what="LEFT — students at gain 128",
        right_what="RIGHT — students at gain 512",
    )
    crops = sorted((edm / "crops").glob("*.png"))
    for i in range(0, len(crops), 4):
        chunk = crops[i:i + 4]
        four_grid(
            prs,
            chunk,
            f"Zoomed teacher/student crops — set {i//4 + 1} (read the gold labels under each)",
            [_humanize_edm_crop(c.stem) for c in chunk],
        )


def _step_label(path: Path) -> str:
    # step_01500.png → "Training step 1500"
    stem = path.stem
    if stem.startswith("step_"):
        try:
            n = int(stem.split("_", 1)[1])
            return f"Training step {n}"
        except ValueError:
            pass
    return stem


def build_live(prs) -> None:
    section_slide(
        prs,
        "Part 6 — Live camera & training progress pictures",
        "What the Pi stream looks like, and how models improve during training (stills; no video file yet)",
    )
    bullets_slide(
        prs,
        "How the live system is supposed to work",
        [
            "Camera on the Pi → denoise on the Hailo chip → video shown in a web page",
            "On-screen numbers show: which chip, resolution, frames per second, lag",
            "We do not have saved .mp4 recordings in this folder yet — only still frames",
            "Training panels below show noisy / cleaned / true-clean as training steps advance",
        ],
    )

    plates = [
        (ROOT / "outputs" / "vis" / "live_fullres_tiled.jpg",
         "Full-size live (or capture) frame after denoise — tiled so large images fit the chip",
         "Shows the real field of view and whether tile seams are visible."),
        (ROOT / "outputs" / "vis" / "live_finetuned_frame.jpg",
         "A live-style frame after finetuning the model on livestream-like content",
         "Finetuned weights applied to a frame that looks like the real stream."),
        (ROOT / "outputs" / "vis" / "1_BEFORE_soft_unaf16.png",
         "OLD result — too soft / oversmoothed (this is what we improved away from)",
         "Baseline “soft” recipe. Detail is lost; looks plastic."),
        (ROOT / "outputs" / "vis" / "2_AFTER_sharp_burst32.png",
         "NEWER result — sharper temporal burst recipe (more detail kept)",
         "Same kind of scene after a sharper training recipe."),
        (ROOT / "outputs" / "vis" / "3_NSA_rebuild_base16_clean.png",
         "Cleaner rebuilt baseline model (base16) — validation picture",
         "NSA rebuild plate for the small base16 network."),
        (ROOT / "outputs" / "vis" / "res_proof.png",
         "Proof picture for resolution / tiling on the live path",
         "Used to check that full resolution still looks correct when split into tiles."),
        (ROOT / "outputs" / "stream_to_gt_panel.png",
         "Livestream frame compared against a clean burst average (quality check)",
         "Left/center/right style panel: live noisy vs model vs ground-truth average."),
        (ROOT / "outputs" / "raw_validation_panel.png",
         "RAW-sensor validation — noisy vs cleaned vs true clean (not phone-JPEG)",
         "Working in the camera’s native Bayer/RAW domain."),
        (ROOT / "outputs" / "high_gain_demo" / "cabinet_H_2_ag512_high_gain.png",
         "Hard night demo — real scene at gain 512 (lots of noise)",
         "Stress test at the highest gain we care about."),
    ]
    for path, title, cap in plates:
        image_plate(prs, path, title, cap, badge="LIVE", what="What this picture is")

    hcg_live = ROOT / "outputs" / "raw_panels_hcg_live"
    step_paths = sorted([p for p in hcg_live.glob("step_*.png")], key=lambda p: p.name)
    if len(step_paths) > 40:
        idxs = sorted(set(
            [0, len(step_paths) - 1]
            + [int(i * (len(step_paths) - 1) / 39) for i in range(40)]
        ))
        step_paths = [step_paths[i] for i in idxs]
    for i in range(0, len(step_paths), 4):
        chunk = step_paths[i:i + 4]
        four_grid(
            prs,
            chunk,
            "HCG / live-like training progress — each tile is noisy | cleaned | true clean",
            [_step_label(c) for c in chunk],
        )
    for p in step_paths[::8]:
        image_plate(
            prs, p,
            f"HCG live-like training snapshot — {_step_label(p).lower()}",
            "Panel from training on dark HCG / livestream-like data. "
            "Watch the middle (cleaned) improve versus the left (noisy).",
            badge="TRAINING",
            what="What this panel is",
        )
    if (hcg_live / "final.png").exists():
        image_plate(
            prs,
            hcg_live / "final.png",
            "Final HCG live-like training panel — last checkpoint of that run",
            "End of the HCG-focused live training series.",
            badge="FINAL",
            what="What this panel is",
        )

    live = ROOT / "outputs" / "panels_live"
    live_steps = sorted(live.glob("step_*.png"), key=lambda p: p.name)
    if len(live_steps) > 24:
        idxs = [int(i * (len(live_steps) - 1) / 23) for i in range(24)]
        live_steps = [live_steps[i] for i in idxs]
    for i in range(0, len(live_steps), 4):
        chunk = live_steps[i:i + 4]
        four_grid(
            prs, chunk,
            "General live curriculum training — noisy | cleaned | true clean",
            [_step_label(c) for c in chunk],
        )
    image_plate(
        prs, live / "final.png",
        "Final panel from the general live training curriculum",
        "Last saved comparison panel for this live training run.",
        badge="FINAL", what="What this panel is",
    )

    raw_live = ROOT / "outputs" / "raw_panels_live"
    raw_steps = sorted(raw_live.glob("step_*.png"), key=lambda p: p.name)
    if len(raw_steps) > 24:
        idxs = [int(i * (len(raw_steps) - 1) / 23) for i in range(24)]
        raw_steps = [raw_steps[i] for i in idxs]
    for i in range(0, len(raw_steps), 4):
        chunk = raw_steps[i:i + 4]
        four_grid(
            prs, chunk,
            "RAW-domain live training progress — noisy | cleaned | true clean",
            [_step_label(c) for c in chunk],
        )
    image_plate(
        prs, raw_live / "final.png",
        "Final RAW-domain live training panel",
        "Last checkpoint comparison while training directly on RAW camera data.",
        badge="FINAL", what="What this panel is",
    )

    for label, rel, plain in [
        ("L1 dark", "outputs/cfm_student_l1_dark_panels",
         "Training with L1 loss on dark scenes"),
        ("L2", "outputs/cfm_student_l2_panels",
         "Training with L2 (mean-squared) loss"),
        ("SWT-rel", "outputs/cfm_student_swtrel_panels",
         "Training with wavelet / relative loss (SWT-rel)"),
        ("SWT-rel2", "outputs/cfm_student_swtrel2_panels",
         "Training with second wavelet / relative loss variant"),
        ("Grad-ratio", "outputs/cfm_gr1e_panels",
         "Training that targets gradient/sharpness ratio"),
    ]:
        d = ROOT / rel
        if not d.exists():
            continue
        panels = sorted(d.glob("step_*.png"))
        for i in range(0, len(panels), 4):
            chunk = panels[i:i + 4]
            four_grid(
                prs, chunk,
                f"{plain} — progress snapshots",
                [_step_label(c) for c in chunk],
            )
        if panels:
            image_plate(
                prs, panels[-1],
                f"Latest snapshot — {plain}",
                f"Last saved panel from {rel}.",
                badge=label.upper(),
                what="What this panel is",
            )


def build_search_and_misc(prs) -> None:
    section_slide(
        prs,
        "Part 7 — Older experiments (architecture search at gain 512)",
        "Historical recipes we tried on the Pi path — useful context, not all are current champions",
    )
    search = ROOT / "outputs" / "pi5_ag512_search"
    imgs = sorted(search.glob("*.png"))
    for i in range(0, len(imgs), 4):
        chunk = imgs[i:i + 4]
        four_grid(
            prs,
            chunk,
            f"Gain-512 recipe search wall — set {i//4 + 1} (labels = experiment names)",
            [c.stem.replace("_", " ") for c in chunk],
        )
    highlights = [
        ("SHOW_current_recipe.png",
         "The recipe we were showing at the time of this search snapshot",
         "Historical “current recipe” panel from the Pi gain-512 search."),
        ("FINAL_unaf16.png",
         "Final small UNAF-16 model result at gain 512",
         "End result for the unaf16 width in that search."),
        ("FINAL_unaf32.png",
         "Final wider UNAF-32 model result at gain 512",
         "End result for the unaf32 width in that search."),
        ("SHARP_trained_unaf16.png",
         "Sharpened / detail-preserving UNAF-16 training result",
         "Variant trained to keep more sharpness."),
        ("TEMPORAL_proof.png",
         "Proof that using several frames in time (temporal) helps",
         "Multi-frame / burst evidence plate."),
        ("RAW_domain_demo.png",
         "Demo of denoising in RAW camera space (not JPEG)",
         "Shows the RAW-domain pipeline visually."),
    ]
    for name, title, cap in highlights:
        image_plate(
            prs, search / name, title, cap,
            badge="OLD SEARCH", what="What this experiment shows",
        )

    section_slide(
        prs,
        "Part 8 — Extra student / teacher panels from loss experiments",
        "Same idea as Part 4: noisy vs cleaned vs true clean, different training recipes",
    )
    misc = [
        ("outputs/cfm_student_panel.png",
         "Default fast student denoise panel",
         "Standard CFM student comparison panel."),
        ("outputs/cfm_student_mixed/cfm_student_panel.png",
         "Student trained on a mixed data recipe",
         "Mixed training set variant."),
        ("outputs/cfm_l2/cfm_student_panel.png",
         "Student trained with L2 loss",
         "L2 loss often looks softer; shown for comparison."),
        ("outputs/cfm_swtrel/cfm_student_panel.png",
         "Student trained with wavelet/relative loss",
         "SWT-rel training variant."),
        ("outputs/cfm_swtrel2/cfm_student_panel.png",
         "Student trained with wavelet/relative loss v2",
         "Second SWT-rel variant."),
        ("outputs/cfm_gr1e/cfm_student_panel.png",
         "Student trained to hit a target sharpness ratio",
         "Gradient-ratio focused training."),
        ("outputs/cfm_tex/cfm_student_panel.png",
         "Student trained with a texture-focused recipe",
         "Texture-oriented student panel."),
        ("outputs/cfm_tex/best_step300.png",
         "Best texture-recipe snapshot at training step 300",
         "Early/best checkpoint from the texture run."),
    ]
    for rel, title, cap in misc:
        image_plate(
            prs, ROOT / rel, title, cap,
            badge="EXPERIMENT", what="What this panel is",
        )

    teacher = ROOT / "outputs" / "cfm_teacher_panels"
    tsteps = [400, 1000, 2000, 3000, 4000, 5000, 6000]
    tpaths = [teacher / f"step_{s:05d}.png" for s in tsteps if (teacher / f"step_{s:05d}.png").exists()]
    for i in range(0, len(tpaths), 4):
        chunk = tpaths[i:i + 4]
        four_grid(
            prs, chunk,
            "Slow teacher model improving over training steps",
            [_step_label(c) for c in chunk],
        )


def build_closing(prs) -> None:
    section_slide(
        prs,
        "Part 9 — Takeaways in plain English",
        "What an outsider should remember after flipping through the pictures",
    )
    bullets_slide(
        prs,
        "What the pictures prove so far",
        [
            "Our fake HCG gain-512 noise now matches the real camera well enough to train on",
            "We built 12,000 training pairs, mostly night-HCG, including live-dark lookalikes",
            "Best small denoisers clean holdout scenes to about 28.4–28.5 dB PSNR",
            "Those models are tiny (~196k weights) and designed for the Hailo chip",
            "Pi CPU alone is too slow for smooth video; Hailo-10H (40 TOPS) is the live path",
            "We still need saved live video clips and on-device FPS overlay screenshots",
        ],
    )
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _textbox(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1),
             "End of the picture gallery", size=36, bold=True, color=FG)
    _textbox(slide, Inches(0.9), Inches(3.5), Inches(11.5), Inches(1.5),
             "Every slide title was written so a new reader can say what the image is "
             "without opening any code.\n\n"
             f"Rebuild:  python scripts/build_portfolio_pptx.py\n"
             f"File:  {OUT.relative_to(ROOT)}",
             size=14, color=MUTED)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building portfolio…")
    build_overview(prs)
    build_compute(prs)
    build_noise(prs)
    build_synth(prs)
    build_champions(prs)
    build_edm(prs)
    build_live(prs)
    build_search_and_misc(prs)
    build_closing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    n = len(prs.slides)
    size_mb = OUT.stat().st_size / (1024 * 1024)
    print(f"Wrote {OUT} · {n} slides · {size_mb:.1f} MB")


if __name__ == "__main__":
    main()
